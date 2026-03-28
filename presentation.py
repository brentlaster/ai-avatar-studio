"""
Presentation Mode — Syncs a PowerPoint deck with a speaker script.

Parses a Markdown script with [SLIDE N] markers, extracts each slide
from a .pptx as an image, generates TTS audio per slide section,
and assembles everything into a single video where each slide is
displayed while the corresponding narration plays.

Script format example:
    [SLIDE 1]
    Welcome everyone! Today we'll be talking about ...

    [SLIDE 2]
    Let's start with the first topic. As you can see on this slide ...

    [SLIDE 3]
    Moving on to the data. These numbers show ...
"""

import os
import re
import subprocess
import shutil
import tempfile
from pathlib import Path
from dataclasses import dataclass
from typing import Optional

from pipeline import (
    generate_speech,
    _generate_speech_coqui_batch,
    AvatarConfig,
    TEMP_DIR,
    OUTPUT_DIR,
    _split_text_into_chunks,
    _concatenate_audio_files,
)


@dataclass
class SlideSegment:
    """One segment of the presentation: a slide number + narration text."""
    slide_number: int
    text: str
    audio_path: Optional[str] = None
    image_path: Optional[str] = None
    video_path: Optional[str] = None


def clean_narration_text(text: str) -> str:
    """
    Clean narration text for speech synthesis:
    - Strip markdown formatting (bold, italic, headers, etc.)
    - Remove stage directions like *[PAUSE]*, *[GESTURE ...]*
    - Remove horizontal rules (---)
    - Remove metadata lines (Duration:, Target pace:, v5 Changes:, etc.)
    - Convert em-dashes to commas for natural TTS pausing
    - Remove link/image markdown
    - Keep the actual spoken words clean and natural
    """
    # Remove code blocks FIRST (before other patterns eat the backticks)
    text = re.sub(r'```[\s\S]*?```', '', text)

    # Remove metadata lines (Duration:, Target pace:, v5 Changes:, etc.)
    text = re.sub(r'^[*]*\s*(?:Duration|Target pace|Total target|v\d+\s+Changes)[*]*\s*:.*$',
                  '', text, flags=re.MULTILINE | re.IGNORECASE)

    # Remove horizontal rules
    text = re.sub(r'^[-*_]{3,}\s*$', '', text, flags=re.MULTILINE)

    # Remove ALL bracketed stage directions (PAUSE, GESTURE, beat, Skip, Wait, etc.)
    # with optional surrounding asterisks/underscores — case insensitive
    text = re.sub(r'[*_]*\[[^\]]*\][*_]*', '', text)

    # Remove HTML tags
    text = re.sub(r'<[^>]+>', '', text)

    # Remove images ![alt](url)
    text = re.sub(r'!\[([^\]]*)\]\([^)]+\)', r'\1', text)

    # Convert links [text](url) to just text
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)

    # Remove reference-style link definitions [label]: url
    text = re.sub(r'^\[[^\]]+\]:\s+.*$', '', text, flags=re.MULTILINE)

    # Remove heading markers (## etc.) but keep the text
    text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)

    # Remove bold/italic markers — run multiple passes to handle nested cases
    # like ***bold italic*** or **bold *nested italic* text**
    for _ in range(3):
        text = re.sub(r'\*{1,3}([^*]*?)\*{1,3}', r'\1', text)
        text = re.sub(r'_{1,3}([^_]*?)_{1,3}', r'\1', text)

    # Remove any remaining lone asterisks or underscores used as emphasis
    text = re.sub(r'(?<!\w)\*+(?!\w)', '', text)
    text = re.sub(r'(?<!\w)_+(?!\w)', '', text)

    # Remove strikethrough
    text = re.sub(r'~~([^~]+)~~', r'\1', text)

    # Remove inline code backticks
    text = re.sub(r'`([^`]*)`', r'\1', text)

    # Remove blockquote markers
    text = re.sub(r'^>\s+', '', text, flags=re.MULTILINE)

    # Remove bullet/list markers (-, *, +, numbered)
    text = re.sub(r'^[\s]*[-*+]\s+', '', text, flags=re.MULTILINE)
    text = re.sub(r'^[\s]*\d+\.\s+', '', text, flags=re.MULTILINE)

    # Remove table formatting
    text = re.sub(r'\|', ' ', text)
    text = re.sub(r'^[\s]*[-:]+[\s]*$', '', text, flags=re.MULTILINE)

    # Replace em-dashes (—) with commas for natural TTS pausing
    # (em-dashes cause TTS to stutter or pause awkwardly)
    text = re.sub(r'\s*—\s*', ', ', text)

    # Replace tildes (used as "approximately") with the word
    text = re.sub(r'~(\d)', r'about \1', text)

    # Replace ellipsis with a period (avoids long dead-air pauses)
    text = re.sub(r'\.{3,}', '.', text)

    # Fix quote boundaries — XTTS generates gibberish when punctuation inside
    # closing quotes butts up against the next sentence without a clear break.
    # e.g.: "Fix the login bug." The model → "Fix the login bug." ... The model
    # Add a newline after closing-quote + punctuation to force a clean sentence break.
    # Remove smart/curly quotes first — TTS handles straight quotes better
    text = text.replace('\u201c', '"').replace('\u201d', '"')
    text = text.replace('\u2018', "'").replace('\u2019', "'")

    # Now force a clean sentence break after closing-quote + punctuation
    # e.g.: "Fix the login bug." The model → two separate sentences
    text = re.sub(r'([.!?])"(\s*)', r'\1"\n\n', text)

    # Clean up punctuation collisions (comma after ? or ! or ., double commas, etc.)
    text = re.sub(r'([?!.])\s*,', r'\1', text)  # "month? ," → "month?"
    text = re.sub(r',\s*([?!.])', r'\1', text)   # ", ." → "."
    text = re.sub(r',\s*,+', ',', text)           # ",," → ","
    text = re.sub(r'([?!.])\s*\1+', r'\1', text)  # ".." → "."

    # Collapse multiple blank lines
    text = re.sub(r'\n{3,}', '\n\n', text)

    # Clean up multiple spaces
    text = re.sub(r'  +', ' ', text)

    # Remove leading/trailing commas on lines
    text = re.sub(r'^\s*,\s*', '', text, flags=re.MULTILINE)
    text = re.sub(r'\s*,\s*$', '', text, flags=re.MULTILINE)

    # Remove leading/trailing whitespace per line
    text = '\n'.join(line.strip() for line in text.split('\n'))

    return text.strip()


def parse_slide_script(script_text: str) -> list[SlideSegment]:
    """
    Parse a script with slide markers into segments.

    Handles two families of marker formats:

    **Bracketed** (original):
      [SLIDE 1]
      [SLIDE 3 — "Title Here"]
      ## [SLIDE 3 — "Title Here"]

    **Heading** (markdown-native):
      ## SLIDE 1: VERSION
      ## SLIDE 3: A SHOW OF HANDS
      ## Slide 10: Wrap-Up

    Text between markers is cleaned of markdown formatting for TTS.
    Text before the first marker is ignored.
    Slides whose body is just "[Skip]" are omitted.
    """
    # Try bracketed format first:  [SLIDE N ...]  or  [SLIDE 11b ...]
    # Capture the number AND an optional letter suffix (e.g. "11b")
    pattern_bracket = r'^\s*#*\s*\[(?:[Ss][Ll][Ii][Dd][Ee])\s+(\d+)([a-zA-Z]?)[^\]]*\]'
    markers = list(re.finditer(pattern_bracket, script_text, re.MULTILINE))

    # If no bracketed markers found, try heading format:  ## SLIDE N: ...
    if not markers:
        pattern_heading = r'^\s*#{1,6}\s+[Ss][Ll][Ii][Dd][Ee]\s+(\d+)([a-zA-Z]?)\b[^\n]*'
        markers = list(re.finditer(pattern_heading, script_text, re.MULTILINE))

    if not markers:
        raise ValueError(
            "No slide markers found in the script.\n"
            "Accepted formats:\n"
            "  [SLIDE 1]\n"
            "  [SLIDE 3 — \"Title Here\"]\n"
            "  [SLIDE 11b — Sub-slide]\n"
            "  ## SLIDE 3: Title Here\n"
        )

    # For the last marker, don't let its text run to end-of-file.
    # Find the first markdown heading that is NOT a SLIDE marker after the
    # last slide — that's where appendix/meta content begins.
    last_marker = markers[-1]
    after_last = last_marker.end()
    # Match headings that don't start with "SLIDE" (case-insensitive)
    end_pattern = re.compile(
        r'^\s*#{1,3}\s+(?![Ss][Ll][Ii][Dd][Ee]\s)(?!\[(?:[Ss][Ll][Ii][Dd][Ee]))',
        re.MULTILINE,
    )
    end_match = end_pattern.search(script_text, after_last)
    if end_match:
        script_boundary = end_match.start()
        print(f"      [parse] Script ends at char {script_boundary:,} "
              f"(of {len(script_text):,}) — meta/appendix content after that ignored")
    else:
        script_boundary = len(script_text)

    # First pass: collect raw slide entries with their script numbers and suffixes
    raw_entries = []
    for i, match in enumerate(markers):
        slide_num = int(match.group(1))
        suffix = match.group(2) if match.lastindex >= 2 else ""  # e.g. "b" from "11b"

        # Text starts after the full marker line
        text_start = match.end()

        # Text ends at the next marker, or at the script boundary
        if i + 1 < len(markers):
            text_end = markers[i + 1].start()
        else:
            text_end = script_boundary

        raw_text = script_text[text_start:text_end].strip()

        # Skip slides explicitly marked [Skip]
        if re.match(r'^\s*\*?\[?\s*[Ss]kip\b.*$', raw_text):
            print(f"      [parse] Slide {slide_num}{suffix}: [Skip] — omitting")
            continue

        raw_entries.append((slide_num, suffix, raw_text))

    # Second pass: compute physical deck slide numbers.
    # Sub-slides like "11b" insert extra slides in the deck after slide 11.
    # Each sub-slide shifts all subsequent physical positions by +1.
    # Example: script 10, 11, 11b, 12 → deck 10, 11, 12, 13
    segments = []
    offset = 0  # cumulative offset from sub-slides
    for slide_num, suffix, raw_text in raw_entries:
        if suffix:
            # Sub-slide: use previous slide's number + 1 in the deck
            offset += 1
            physical_num = slide_num + offset
            print(f"      [parse] Slide {slide_num}{suffix} → deck slide {physical_num} (sub-slide)")
        else:
            physical_num = slide_num + offset

        clean_text = clean_narration_text(raw_text)
        if clean_text:
            segments.append(SlideSegment(slide_number=physical_num, text=clean_text))

    return segments


def _find_soffice() -> str | None:
    """Find the LibreOffice 'soffice' binary, including macOS app bundles."""
    import shutil

    # Check if it's on PATH (Linux / Homebrew installs)
    path = shutil.which("soffice")
    if path:
        return path

    # macOS: LibreOffice installs into /Applications
    mac_paths = [
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        os.path.expanduser("~/Applications/LibreOffice.app/Contents/MacOS/soffice"),
    ]
    for p in mac_paths:
        if os.path.isfile(p):
            return p

    return None


def _find_pdftoppm() -> str | None:
    """Find the pdftoppm binary (part of poppler)."""
    import shutil

    path = shutil.which("pdftoppm")
    if path:
        return path

    # Homebrew on Apple Silicon / Intel
    for prefix in ["/opt/homebrew/bin/pdftoppm", "/usr/local/bin/pdftoppm"]:
        if os.path.isfile(prefix):
            return prefix

    return None


def extract_slides_as_images(
    pptx_path: str,
    output_dir: str,
    dpi: int = 150,
) -> dict[int, str]:
    """
    Convert a .pptx file to individual slide images.
    Returns a dict mapping slide number (1-based) to image path.

    Tries LibreOffice + pdftoppm first; falls back to python-pptx + Pillow
    if those tools aren't installed.
    """
    pptx_path = Path(pptx_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    soffice = _find_soffice()
    pdftoppm = _find_pdftoppm()

    if soffice and pdftoppm:
        return _extract_slides_libreoffice(pptx_path, output_dir, dpi, soffice, pdftoppm)
    else:
        missing = []
        if not soffice:
            missing.append("LibreOffice (soffice)")
        if not pdftoppm:
            missing.append("pdftoppm (poppler)")
        print(f"      Note: {', '.join(missing)} not found — using python-pptx fallback")
        print(f"      For best quality, install: brew install --cask libreoffice && brew install poppler")
        return _extract_slides_pptx_fallback(pptx_path, output_dir)


def _extract_slides_libreoffice(
    pptx_path: Path,
    output_dir: Path,
    dpi: int,
    soffice: str,
    pdftoppm: str,
) -> dict[int, str]:
    """Extract slides using LibreOffice → PDF → pdftoppm (highest quality)."""

    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_path = Path(tmp_dir)

        # Step 1: Convert PPTX to PDF via LibreOffice
        print("      Converting slides to PDF via LibreOffice ...")
        env = os.environ.copy()
        env["SAL_USE_VCLPLUGIN"] = "svp"

        result = subprocess.run(
            [
                soffice, "--headless",
                "--convert-to", "pdf",
                "--outdir", str(tmp_path),
                str(pptx_path),
            ],
            capture_output=True, text=True,
            env=env,
            timeout=120,
        )

        pdf_path = tmp_path / f"{pptx_path.stem}.pdf"
        if result.returncode != 0 or not pdf_path.exists():
            raise RuntimeError(
                f"LibreOffice PDF conversion failed:\n{result.stderr}"
            )

        # Step 2: Render each PDF page as a JPEG image
        print("      Rendering slide images ...")
        result = subprocess.run(
            [
                pdftoppm,
                "-jpeg", "-r", str(dpi),
                str(pdf_path),
                str(output_dir / "slide"),
            ],
            capture_output=True, text=True,
            timeout=120,
        )

        if result.returncode != 0:
            raise RuntimeError(
                f"pdftoppm image conversion failed:\n{result.stderr}"
            )

    # Collect the generated images (named slide-01.jpg, slide-02.jpg, etc.)
    slide_images = sorted(output_dir.glob("slide-*.jpg"))

    if not slide_images:
        raise RuntimeError("No slide images were generated from the PPTX.")

    # Map 1-based slide number to image path
    slide_map = {}
    for idx, img_path in enumerate(slide_images, start=1):
        slide_map[idx] = str(img_path)

    print(f"      Extracted {len(slide_map)} slide images (LibreOffice)")
    return slide_map


def _extract_slides_pptx_fallback(
    pptx_path: Path,
    output_dir: Path,
) -> dict[int, str]:
    """
    Fallback: render slides using python-pptx + Pillow.
    Produces simpler renders (text + solid backgrounds) but works
    without any external dependencies.
    """
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches, Pt, Emu
    from PIL import Image, ImageDraw, ImageFont

    prs = PptxPresentation(str(pptx_path))
    slide_width = prs.slide_width or Emu(12192000)   # default 10"
    slide_height = prs.slide_height or Emu(6858000)  # default 7.5"

    # Render at 1920px wide, scale height proportionally
    img_w = 1920
    img_h = int(img_w * (slide_height / slide_width))

    # Try to find a reasonable font
    font_paths = [
        "/System/Library/Fonts/Helvetica.ttc",
        "/System/Library/Fonts/SFNSText.ttf",
        "/Library/Fonts/Arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    font_path = None
    for fp in font_paths:
        if os.path.isfile(fp):
            font_path = fp
            break

    def get_font(size):
        try:
            if font_path:
                return ImageFont.truetype(font_path, size)
        except Exception:
            pass
        return ImageFont.load_default()

    slide_map = {}
    for slide_num, slide in enumerate(prs.slides, start=1):
        img = Image.new("RGB", (img_w, img_h), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)

        # Extract text from all shapes and render
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)

        if texts:
            # Render title larger, body smaller
            y_pos = int(img_h * 0.1)
            title_font = get_font(48)
            body_font = get_font(32)

            for i, text in enumerate(texts):
                font = title_font if i == 0 else body_font
                # Word wrap
                words = text.split()
                lines = []
                current_line = ""
                for word in words:
                    test = f"{current_line} {word}".strip()
                    bbox = draw.textbbox((0, 0), test, font=font)
                    if bbox[2] < img_w - 160:
                        current_line = test
                    else:
                        if current_line:
                            lines.append(current_line)
                        current_line = word
                if current_line:
                    lines.append(current_line)

                for line in lines:
                    bbox = draw.textbbox((0, 0), line, font=font)
                    text_w = bbox[2] - bbox[0]
                    x = (img_w - text_w) // 2
                    draw.text((x, y_pos), line, fill=(33, 33, 33), font=font)
                    y_pos += int((bbox[3] - bbox[1]) * 1.4)

                y_pos += 20  # gap between shapes
        else:
            # Blank slide — just add slide number
            font = get_font(36)
            draw.text((img_w // 2 - 40, img_h // 2), f"Slide {slide_num}", fill=(150, 150, 150), font=font)

        out_path = output_dir / f"slide-{slide_num:02d}.jpg"
        img.save(str(out_path), "JPEG", quality=90)
        slide_map[slide_num] = str(out_path)

    print(f"      Extracted {len(slide_map)} slide images (python-pptx fallback)")
    print(f"      Tip: Install LibreOffice for higher quality renders:")
    print(f"           brew install --cask libreoffice && brew install poppler")
    return slide_map


def create_slide_video(
    image_path: str,
    audio_path: str,
    output_path: str,
    lead_in: float = 0.3,
) -> str:
    """
    Create a video segment: static slide image + audio narration.

    A short silent lead-in (default 0.3s) is prepended so the slide
    appears on screen *before* the narration begins.  This prevents the
    common perception that the slide is "late" relative to the voice.
    Reduced from 0.5s because XTTS chunks already have natural
    onset timing and leading silence is now trimmed.
    """
    # Get audio sample rate so the silence matches exactly
    probe_rate_cmd = [
        "ffprobe", "-v", "quiet",
        "-select_streams", "a:0",
        "-show_entries", "stream=sample_rate",
        "-of", "csv=p=0",
        audio_path,
    ]
    rate_result = subprocess.run(probe_rate_cmd, capture_output=True, text=True)
    try:
        sample_rate = int(rate_result.stdout.strip())
    except (ValueError, AttributeError):
        sample_rate = 24000

    # Build audio with lead-in silence prepended using ffmpeg's adelay filter.
    # adelay takes milliseconds; we also add a tiny 150ms tail of silence so the
    # last word doesn't feel clipped. A 50ms fade-out at the very end eliminates
    # stray XTTS speech artifacts at slide boundaries.
    lead_ms = int(lead_in * 1000)
    padded_audio_path = os.path.splitext(audio_path)[0] + "_padded.wav"

    # Step 1: create padded audio (lead-in silence + 150ms tail silence)
    pad_cmd_nofade = [
        "ffmpeg", "-y",
        "-i", audio_path,
        "-af", f"adelay={lead_ms}|{lead_ms},apad=pad_dur=0.15",
        "-ar", str(sample_rate),
        padded_audio_path,
    ]
    pad_result = subprocess.run(pad_cmd_nofade, capture_output=True, text=True, timeout=120)
    if pad_result.returncode != 0:
        # Fallback: use original audio without padding
        print(f"      Warning: Could not add lead-in silence, using original audio")
        padded_audio_path = audio_path

    # Get padded audio duration for the video
    probe_cmd = [
        "ffprobe", "-v", "quiet",
        "-show_entries", "format=duration",
        "-of", "csv=p=0",
        padded_audio_path,
    ]
    result = subprocess.run(probe_cmd, capture_output=True, text=True)
    try:
        duration = float(result.stdout.strip())
    except ValueError:
        duration = 10.0

    # Step 2: apply fade-out to the last 80ms to kill stray XTTS artifacts
    fade_start = max(0.0, duration - 0.08)
    faded_audio_path = os.path.splitext(audio_path)[0] + "_faded.wav"
    fade_cmd = [
        "ffmpeg", "-y",
        "-i", padded_audio_path,
        "-af", f"afade=t=out:st={fade_start:.3f}:d=0.08",
        "-ar", str(sample_rate),
        faded_audio_path,
    ]
    fade_result = subprocess.run(fade_cmd, capture_output=True, text=True, timeout=120)
    if fade_result.returncode == 0:
        final_audio = faded_audio_path
    else:
        final_audio = padded_audio_path

    # Create video: static image looped for EXACTLY the audio duration + audio track.
    # Using explicit -t duration instead of -shortest prevents A/V duration mismatch
    # that causes cumulative timing drift when segments are concatenated.
    cmd = [
        "ffmpeg", "-y",
        "-loop", "1",
        "-i", image_path,
        "-i", final_audio,
        "-c:v", "libx264",
        "-tune", "stillimage",
        "-c:a", "aac", "-b:a", "192k",
        "-vf", "scale=1920:1080:force_original_aspect_ratio=decrease,pad=1920:1080:(ow-iw)/2:(oh-ih)/2:black",
        "-pix_fmt", "yuv420p",
        "-t", f"{duration:.3f}",
        "-movflags", "+faststart",
        output_path,
    ]

    result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
    if result.returncode != 0:
        raise RuntimeError(f"Slide video creation failed:\n{result.stderr}")

    # Clean up temp audio files
    for tmp in [padded_audio_path, faded_audio_path]:
        if tmp != audio_path:
            try:
                os.remove(tmp)
            except OSError:
                pass

    return output_path


def concatenate_videos(video_paths: list[str], output_path: str) -> str:
    """Concatenate multiple video segments into one final presentation video.

    Always re-encodes during concatenation to guarantee proper A/V timestamp
    alignment.  The old ``-c copy`` approach was faster but caused cumulative
    timing drift: each segment's audio (AAC, with priming samples) and video
    (25 fps still-image loop) had slightly different durations, and the small
    errors compounded over many slides—producing visible slide-lag.
    """
    if len(video_paths) == 1:
        shutil.copy2(video_paths[0], output_path)
        return output_path

    # Create concat list file
    list_path = os.path.join(TEMP_DIR, "video_concat_list.txt")
    with open(list_path, "w") as f:
        for path in video_paths:
            escaped = path.replace("'", "'\\''")
            f.write(f"file '{escaped}'\n")

    # Re-encode during concat to force correct timestamp alignment.
    # Uses CRF 20 (visually lossless for still-image slides) and fast preset
    # since the video is just static slides—encoding is quick.
    cmd = [
        "ffmpeg", "-y",
        "-f", "concat", "-safe", "0",
        "-i", list_path,
        "-c:v", "libx264", "-crf", "20", "-preset", "fast",
        "-tune", "stillimage",
        "-c:a", "aac", "-b:a", "192k",
        "-pix_fmt", "yuv420p",
        "-movflags", "+faststart",
        output_path,
    ]

    result = subprocess.run(cmd, capture_output=True, text=True, timeout=600)
    if result.returncode != 0:
        raise RuntimeError(f"Video concatenation failed:\n{result.stderr}")

    return output_path


def _extract_slide_notes(pptx_path: str) -> dict[int, str]:
    """
    Extract speaker notes from each slide in a .pptx file.
    Returns a dict mapping slide number (1-based) to notes text.
    Slides without notes return an empty string.
    """
    notes_map: dict[int, str] = {}
    try:
        from pptx import Presentation as PptxPresentation
        prs = PptxPresentation(pptx_path)
        for i, slide in enumerate(prs.slides, start=1):
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    notes_map[i] = notes_text
        if notes_map:
            print(f"      Extracted speaker notes from {len(notes_map)} slides")
    except Exception as e:
        print(f"      Warning: Could not extract slide notes: {e}")
    return notes_map


def generate_presentation(
    pptx_path: str,
    script_text: str,
    voice_id: str,
    config: AvatarConfig = AvatarConfig(),
    output_path: Optional[str] = None,
    progress_callback=None,
    start_slide: int = 0,
    end_slide: int = 0,
) -> str:
    """
    Full presentation generation pipeline:
      1. Parse script for [SLIDE N] markers
      2. Extract slides from PPTX as images
      3. Generate TTS audio for each slide's narration
      4. Create a video per slide (image + audio)
      5. Concatenate all segments into the final video

    start_slide / end_slide: if non-zero, only generate slides in that
    range (inclusive). 0 means "no limit" on that end.

    Returns path to the final presentation video.
    """
    if output_path is None:
        output_path = os.path.join(OUTPUT_DIR, "presentation.mp4")

    # Clean and create working directories (remove stale files from previous runs)
    slides_dir = os.path.join(TEMP_DIR, "presentation_slides")
    segments_dir = os.path.join(TEMP_DIR, "presentation_segments")
    for d in (slides_dir, segments_dir):
        if os.path.exists(d):
            shutil.rmtree(d)
    os.makedirs(slides_dir, exist_ok=True)
    os.makedirs(segments_dir, exist_ok=True)

    print("=" * 60)
    print("  Presentation Mode")
    print("=" * 60)

    # Step 1: Parse the script
    print("\n[1/4] Parsing slide script ...")
    segments = parse_slide_script(script_text)
    print(f"      Found {len(segments)} slide segments")

    # Apply slide range filter
    if start_slide > 0 or end_slide > 0:
        original_count = len(segments)
        if start_slide > 0:
            segments = [s for s in segments if s.slide_number >= start_slide]
        if end_slide > 0:
            segments = [s for s in segments if s.slide_number <= end_slide]
        range_desc = f"slides {start_slide or 'start'}-{end_slide or 'end'}"
        print(f"      Filtered to {range_desc}: {original_count} → {len(segments)} segments")
        if not segments:
            raise ValueError(
                f"No slides in range {start_slide}-{end_slide}. "
                f"Available slides: {', '.join(str(s.slide_number) for s in parse_slide_script(script_text))}"
            )

    for seg in segments:
        preview = seg.text[:60].replace('\n', ' ')
        print(f"        Slide {seg.slide_number}: \"{preview}...\"")

    # Step 2: Extract slides as images
    print(f"\n[2/4] Extracting slides from {Path(pptx_path).name} ...")
    slide_images = extract_slides_as_images(pptx_path, slides_dir)

    # Validate that all referenced slides exist
    max_slide = max(slide_images.keys()) if slide_images else 0
    for seg in segments:
        if seg.slide_number not in slide_images:
            if seg.slide_number > max_slide:
                raise ValueError(
                    f"Script references [SLIDE {seg.slide_number}] but the deck "
                    f"only has {max_slide} slides."
                )
            raise ValueError(
                f"Slide {seg.slide_number} not found in extracted images."
            )
        seg.image_path = slide_images[seg.slide_number]

    # Step 3: Generate audio for each segment
    print(f"\n[3/4] Generating narration audio ({len(segments)} segments) ...")
    total_chars = sum(len(seg.text) for seg in segments)
    print(f"      Total script: {total_chars:,} characters")

    if config.tts_engine == "coqui_xtts" and len(segments) > 1:
        # --- BATCH MODE: load Coqui model once for all slides ---
        print(f"      Using batch mode (model loads once for all {len(segments)} slides)")
        batch_items = []
        # Use sequential index for file names to avoid collisions when
        # sub-slides (e.g. 11b) produce duplicate slide numbers.
        for idx, seg in enumerate(segments):
            audio_path = os.path.join(segments_dir, f"narration_seg{idx:03d}_slide{seg.slide_number:03d}.wav")
            batch_items.append({
                "text": seg.text,
                "output_path": audio_path,
                "label": f"Slide {seg.slide_number}",
            })

        if progress_callback:
            progress_callback(0.25, desc="Generating all slide audio (batch mode)...")

        _generate_speech_coqui_batch(batch_items, config)

        # Map outputs back to segments
        for idx, seg in enumerate(segments):
            audio_path = os.path.join(segments_dir, f"narration_seg{idx:03d}_slide{seg.slide_number:03d}.wav")
            if os.path.exists(audio_path):
                seg.audio_path = audio_path
            else:
                raise RuntimeError(f"Batch audio not found for slide {seg.slide_number}: {audio_path}")

        if progress_callback:
            progress_callback(0.7, desc="All slide audio generated")
    else:
        # --- SEQUENTIAL MODE: one subprocess per segment (ElevenLabs or single slide) ---
        for i, seg in enumerate(segments):
            print(f"\n      Slide {seg.slide_number} ({i + 1}/{len(segments)}, {len(seg.text):,} chars) ...")
            preview_text = seg.text[:200].replace('\n', ' ↵ ')
            print(f"      Narration: \"{preview_text}{'...' if len(seg.text) > 200 else ''}\"")
            audio_ext = ".wav" if config.tts_engine == "coqui_xtts" else ".mp3"
            audio_path = os.path.join(segments_dir, f"narration_seg{i:03d}_slide{seg.slide_number:03d}{audio_ext}")
            generate_speech(
                script_text=seg.text,
                voice_id=voice_id,
                config=config,
                output_path=audio_path,
            )
            wav_alt = os.path.splitext(audio_path)[0] + ".wav"
            mp3_alt = os.path.splitext(audio_path)[0] + ".mp3"
            if os.path.exists(audio_path):
                seg.audio_path = audio_path
            elif os.path.exists(wav_alt):
                seg.audio_path = wav_alt
            elif os.path.exists(mp3_alt):
                seg.audio_path = mp3_alt
            else:
                raise RuntimeError(f"Audio file not found for slide {seg.slide_number}: tried {audio_path}")

            if progress_callback:
                pct = 0.2 + 0.5 * ((i + 1) / len(segments))
                progress_callback(pct, desc=f"Generated audio for slide {seg.slide_number}")

    # Step 4: Create video segments (parallel) and concatenate
    print(f"\n[4/4] Assembling presentation video ...")
    from concurrent.futures import ThreadPoolExecutor, as_completed

    def _make_slide_video(idx_seg):
        """Worker function for parallel video assembly."""
        idx, seg = idx_seg
        seg_video_path = os.path.join(segments_dir, f"segment_seg{idx:03d}_slide{seg.slide_number:03d}.mp4")
        create_slide_video(seg.image_path, seg.audio_path, seg_video_path)
        seg.video_path = seg_video_path
        return idx, seg_video_path

    video_segments = []
    # Use up to 4 parallel ffmpeg workers (CPU-bound, not memory-heavy)
    max_workers = min(4, len(segments))
    if max_workers > 1:
        print(f"      Creating {len(segments)} slide videos ({max_workers} parallel workers) ...")
        video_map = {}
        with ThreadPoolExecutor(max_workers=max_workers) as pool:
            futures = {pool.submit(_make_slide_video, (idx, seg)): idx for idx, seg in enumerate(segments)}
            for future in as_completed(futures):
                idx, video_path = future.result()
                video_map[idx] = video_path
                print(f"      Segment {idx} (slide {segments[idx].slide_number}) video ready")

        # Maintain original order using sequential index
        for idx in range(len(segments)):
            video_segments.append(video_map[idx])
    else:
        for idx, seg in enumerate(segments):
            print(f"      Creating video for slide {seg.slide_number} ...")
            seg_video_path = os.path.join(segments_dir, f"segment_seg{idx:03d}_slide{seg.slide_number:03d}.mp4")
            create_slide_video(seg.image_path, seg.audio_path, seg_video_path)
            seg.video_path = seg_video_path
            video_segments.append(seg_video_path)

    if progress_callback:
        progress_callback(0.95, desc="Concatenating final video...")

    print(f"\n      Concatenating {len(video_segments)} segments ...")
    concatenate_videos(video_segments, output_path)

    # Get duration of each segment's audio for timeline
    seg_durations = []
    for seg in segments:
        probe_cmd = [
            "ffprobe", "-v", "quiet",
            "-show_entries", "format=duration",
            "-of", "csv=p=0",
            seg.audio_path,
        ]
        result = subprocess.run(probe_cmd, capture_output=True, text=True)
        try:
            dur = float(result.stdout.strip())
        except (ValueError, AttributeError):
            dur = 5.0  # fallback
        seg_durations.append(dur)

    # Calculate total duration
    total_duration = sum(seg_durations)
    minutes = int(total_duration // 60)
    seconds = int(total_duration % 60)
    print(f"\n      Total presentation duration: {minutes}m {seconds}s")

    # Extract slide notes from the PPTX for display in the viewer
    slide_notes = _extract_slide_notes(pptx_path)

    # Build timeline data for the in-app script viewer
    timeline = []
    cumulative = 0.0
    for i, seg in enumerate(segments):
        notes = slide_notes.get(seg.slide_number, "")
        timeline.append({
            "slide": seg.slide_number,
            "start": round(cumulative, 2),
            "end": round(cumulative + seg_durations[i], 2),
            "text": seg.text,
            "notes": notes,
        })
        cumulative += seg_durations[i]

    # Store timeline as JSON alongside the video for the Gradio app to read
    timeline_path = os.path.splitext(output_path)[0] + "_timeline.json"
    try:
        import json
        with open(timeline_path, "w", encoding="utf-8") as f:
            json.dump(timeline, f, indent=2, ensure_ascii=False)
        print(f"      Script timeline saved to {timeline_path}")
    except Exception as e:
        print(f"      Warning: Could not save script timeline: {e}")

    # Generate standalone HTML viewer (video + synced script in one file)
    try:
        viewer_path = generate_standalone_viewer(output_path, timeline)
        if viewer_path:
            print(f"      Standalone viewer: {viewer_path}")
    except Exception as e:
        print(f"      Warning: Could not generate standalone viewer: {e}")

    # Generate mobile-friendly viewer (small HTML + separate MP4)
    try:
        mobile_path = generate_mobile_viewer(output_path, timeline)
        if mobile_path:
            print(f"      Mobile viewer:     {mobile_path}")
    except Exception as e:
        print(f"      Warning: Could not generate mobile viewer: {e}")

    print(f"      Presentation saved to {output_path}")
    print("=" * 60)
    print("  Done!")
    print("=" * 60)

    return output_path


def build_script_viewer_html(timeline: list) -> str:
    """
    Build HTML for an in-app script viewer panel.
    Shows each slide's narration text with timestamps in a scrollable,
    styled panel that can be embedded in the Gradio UI.
    """
    import html as html_mod

    if not timeline:
        return ""

    total_duration = timeline[-1]["end"] if timeline else 0
    total_min = int(total_duration // 60)
    total_sec = int(total_duration % 60)

    seg_blocks = ""
    for i, t in enumerate(timeline):
        escaped_text = html_mod.escape(t["text"]).replace("\n", "<br>")
        start_min = int(t["start"] // 60)
        start_sec = int(t["start"] % 60)
        end_min = int(t["end"] // 60)
        end_sec = int(t["end"] % 60)
        duration = t["end"] - t["start"]

        notes_html = ""
        notes_text = t.get("notes", "")
        if notes_text:
            escaped_notes = html_mod.escape(notes_text).replace("\n", "<br>")
            notes_html = (
                f'<div class="sv-notes">'
                f'<span class="sv-notes-label">Slide Notes</span>'
                f'{escaped_notes}</div>'
            )

        seg_blocks += f'''
        <div class="sv-seg" id="sv-seg-{i}">
            <div class="sv-header">
                <span class="sv-slide">Slide {t['slide']}</span>
                <span class="sv-time">{start_min}:{start_sec:02d} — {end_min}:{end_sec:02d} ({duration:.0f}s)</span>
            </div>
            <div class="sv-text">{escaped_text}</div>
            {notes_html}
        </div>
'''

    html = f'''
<div class="script-viewer-container">
    <div class="sv-title">Speaker Script &nbsp;·&nbsp; {len(timeline)} slides &nbsp;·&nbsp; {total_min}m {total_sec}s total</div>
    <div class="sv-scroll">
{seg_blocks}
    </div>
</div>
<style>
.script-viewer-container {{
    font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
    background: #f8fafc;
    border-radius: 10px;
    overflow: hidden;
    border: 1px solid #e2e8f0;
}}
.sv-title {{
    padding: 12px 18px;
    font-size: 13px;
    font-weight: 600;
    color: #475569;
    text-transform: uppercase;
    letter-spacing: 0.5px;
    background: #f1f5f9;
    border-bottom: 1px solid #e2e8f0;
}}
.sv-scroll {{
    max-height: 480px;
    overflow-y: auto;
    padding: 12px;
}}
.sv-seg {{
    padding: 14px 16px;
    margin-bottom: 8px;
    border-radius: 8px;
    border-left: 3px solid #cbd5e1;
    background: #ffffff;
    transition: all 0.2s ease;
    box-shadow: 0 1px 2px rgba(0,0,0,0.04);
}}
.sv-seg:hover {{
    background: #f1f5f9;
    border-left-color: #3b82f6;
}}
.sv-header {{
    display: flex;
    justify-content: space-between;
    align-items: center;
    margin-bottom: 8px;
}}
.sv-slide {{
    font-weight: 600;
    font-size: 13px;
    color: #2563eb;
}}
.sv-time {{
    font-size: 12px;
    color: #94a3b8;
    font-variant-numeric: tabular-nums;
}}
.sv-text {{
    font-size: 14px;
    line-height: 1.65;
    color: #1e293b;
}}
.sv-notes {{
    margin-top: 10px;
    padding: 10px 12px;
    background: #fefce8;
    border: 1px solid #fde68a;
    border-radius: 6px;
    font-size: 13px;
    line-height: 1.55;
    color: #713f12;
}}
.sv-notes-label {{
    display: block;
    font-size: 11px;
    font-weight: 600;
    text-transform: uppercase;
    letter-spacing: 0.5px;
    color: #a16207;
    margin-bottom: 4px;
}}
.sv-scroll::-webkit-scrollbar {{
    width: 6px;
}}
.sv-scroll::-webkit-scrollbar-track {{
    background: transparent;
}}
.sv-scroll::-webkit-scrollbar-thumb {{
    background: #cbd5e1;
    border-radius: 3px;
}}
</style>
'''
    return html


def generate_standalone_viewer(video_path: str, timeline: list) -> str:
    """
    Generate a self-contained HTML file with an embedded video player
    and synced scrolling script. The video is base64-encoded into the
    HTML so the file works when opened directly in any browser
    (no server needed, no 'about:blank#blocked' issues).

    Returns the path to the generated HTML file.
    """
    import html as html_mod
    import base64

    if not timeline or not os.path.exists(video_path):
        return ""

    viewer_path = os.path.splitext(video_path)[0] + "_viewer.html"

    # Base64-encode the video so the HTML is fully self-contained
    print("      Encoding video for standalone viewer (this may take a moment) ...")
    with open(video_path, "rb") as vf:
        video_b64 = base64.b64encode(vf.read()).decode("ascii")

    total_duration = timeline[-1]["end"] if timeline else 0
    total_min = int(total_duration // 60)
    total_sec = int(total_duration % 60)

    # Build segment HTML blocks and collect notes for the JS notes data array
    seg_blocks = ""
    notes_js_entries = []
    for i, t in enumerate(timeline):
        escaped_text = html_mod.escape(t["text"]).replace("\n", "<br>")
        start_min = int(t["start"] // 60)
        start_sec = int(t["start"] % 60)
        notes_html = ""
        notes_text = t.get("notes", "")
        if notes_text:
            escaped_notes = html_mod.escape(notes_text).replace("\n", "<br>")
            notes_html = (
                f'<div class="seg-notes">'
                f'<span class="seg-notes-label">Slide Notes</span>'
                f'{escaped_notes}</div>'
            )
        # Escape notes for JS string (handle quotes and newlines)
        js_notes = html_mod.escape(notes_text).replace("\\", "\\\\").replace('"', '\\"').replace("\n", "<br>")
        notes_js_entries.append(f'{{start:{t["start"]},end:{t["end"]},slide:{t["slide"]},notes:"{js_notes}"}}')
        seg_blocks += f'''
        <div class="seg" id="seg-{i}" data-start="{t['start']}" data-end="{t['end']}">
            <div class="seg-header">
                Slide {t['slide']}
                <span class="seg-time">{start_min}:{start_sec:02d}</span>
            </div>
            <div class="seg-text">{escaped_text}</div>
            {notes_html}
        </div>'''
    notes_js_array = ",".join(notes_js_entries)
    has_any_notes = any(t.get("notes", "") for t in timeline)

    # Pre-build the notes bar HTML (avoids nested quotes in f-string)
    if has_any_notes:
        notes_bar_html = (
            '<div class="notes-bar empty" id="notesBar">'
            '<div class="notes-bar-label">Slide Notes</div>'
            '<div class="notes-bar-text placeholder" id="notesText">'
            'Notes will appear here as the presentation plays</div></div>'
        )
    else:
        notes_bar_html = ""

    html_content = f'''<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>Presentation Viewer — {len(timeline)} slides, {total_min}m {total_sec}s</title>
<style>
* {{ margin: 0; padding: 0; box-sizing: border-box; }}
body {{ font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
       background: #1a1a2e; color: #e0e0e0; height: 100vh; overflow: hidden; }}
.container {{ display: flex; height: 100vh; }}
.video-panel {{ flex: 1; display: flex; flex-direction: column; align-items: center;
               justify-content: center; padding: 20px; background: #16213e; }}
.video-panel video {{ max-width: 100%; max-height: 70vh; border-radius: 8px;
                      box-shadow: 0 4px 20px rgba(0,0,0,0.5); }}
.video-panel h2 {{ color: #a0c4ff; margin-bottom: 12px; font-size: 14px;
                   letter-spacing: 1px; text-transform: uppercase; }}
.speed-bar {{ display: flex; gap: 6px; margin-top: 14px; align-items: center; }}
.speed-bar span {{ font-size: 12px; color: #7a8ba8; margin-right: 4px; }}
.speed-btn {{ background: rgba(255,255,255,0.08); border: 1px solid #334155;
              color: #94a3b8; padding: 5px 12px; border-radius: 6px; cursor: pointer;
              font-size: 13px; transition: all 0.15s; }}
.speed-btn:hover {{ background: rgba(255,255,255,0.15); color: #e0e0e0; }}
.speed-btn.active {{ background: #539cff; color: #fff; border-color: #539cff; }}
.script-panel {{ width: 420px; min-width: 350px; overflow-y: auto; padding: 20px;
                background: #0f3460; border-left: 2px solid #1a1a2e; }}
.script-panel h2 {{ color: #a0c4ff; margin-bottom: 16px; font-size: 14px;
                    letter-spacing: 1px; text-transform: uppercase;
                    position: sticky; top: 0; background: #0f3460; padding: 8px 0; z-index: 1; }}
.seg {{ padding: 14px 16px; margin-bottom: 10px; border-radius: 8px;
        border-left: 3px solid transparent; transition: all 0.3s ease;
        cursor: pointer; background: rgba(255,255,255,0.03); }}
.seg:hover {{ background: rgba(255,255,255,0.08); }}
.seg.active {{ background: rgba(83,156,255,0.15); border-left-color: #539cff; }}
.seg-header {{ font-weight: 600; color: #539cff; margin-bottom: 6px;
              display: flex; justify-content: space-between; align-items: center; font-size: 13px; }}
.seg-time {{ font-weight: 400; color: #7a8ba8; font-size: 12px; }}
.seg-text {{ font-size: 14px; line-height: 1.6; color: #c8d6e5; }}
.seg.active .seg-text {{ color: #f0f0f0; }}
.seg-notes {{ margin-top: 10px; padding: 10px 12px; background: rgba(251,191,36,0.12);
              border: 1px solid rgba(251,191,36,0.3); border-radius: 6px;
              font-size: 13px; line-height: 1.55; color: #fde68a; }}
.seg-notes-label {{ display: block; font-size: 11px; font-weight: 600;
                    text-transform: uppercase; letter-spacing: 0.5px;
                    color: #fbbf24; margin-bottom: 4px; }}
.notes-bar {{ width: 100%; margin-top: 14px; padding: 12px 16px;
             background: rgba(251,191,36,0.10); border: 1px solid rgba(251,191,36,0.25);
             border-radius: 8px; min-height: 48px; max-height: 120px; overflow-y: auto;
             transition: opacity 0.3s ease; }}
.notes-bar.empty {{ opacity: 0.4; }}
.notes-bar-label {{ font-size: 11px; font-weight: 600; text-transform: uppercase;
                    letter-spacing: 0.5px; color: #fbbf24; margin-bottom: 4px; }}
.notes-bar-text {{ font-size: 13px; line-height: 1.5; color: #fde68a; }}
.notes-bar-text.placeholder {{ color: #7a8ba8; font-style: italic; }}
.script-panel::-webkit-scrollbar {{ width: 6px; }}
.script-panel::-webkit-scrollbar-track {{ background: transparent; }}
.script-panel::-webkit-scrollbar-thumb {{ background: #475569; border-radius: 3px; }}
@media (max-width: 768px) {{
    body {{ height: auto; overflow-y: auto; }}
    .container {{ flex-direction: column; height: auto; overflow-y: visible; }}
    .video-panel {{ padding: 12px; flex: none; }}
    .video-panel video {{ max-height: 35vh; width: 100%; }}
    .script-panel {{ width: 100%; min-width: unset; max-height: none;
                     overflow-y: visible;
                     border-left: none; border-top: 2px solid #1a1a2e; }}
    .speed-bar {{ flex-wrap: wrap; }}
    .notes-bar {{ max-height: 100px; }}
}}
</style>
</head>
<body>
<div class="container">
    <div class="video-panel">
        <h2>Presentation</h2>
        <video id="vid" controls playsinline webkit-playsinline></video>
        <div id="loadingMsg" style="color:#7a8ba8;font-size:13px;margin-top:8px;">Loading video...</div>
        <div class="speed-bar">
            <span>Speed:</span>
            <button class="speed-btn" data-speed="0.5">0.5x</button>
            <button class="speed-btn" data-speed="0.75">0.75x</button>
            <button class="speed-btn active" data-speed="1">1x</button>
            <button class="speed-btn" data-speed="1.25">1.25x</button>
            <button class="speed-btn" data-speed="1.5">1.5x</button>
            <button class="speed-btn" data-speed="2">2x</button>
        </div>
        {notes_bar_html}
    </div>
    <div class="script-panel" id="scriptPanel">
        <h2>Speaker Script</h2>
{seg_blocks}
    </div>
</div>
<script>
// Convert base64 video to Blob URL for mobile compatibility
// (mobile Safari does not support data: URIs on <video> elements)
const videoB64 = "{video_b64}";
const loadMsg = document.getElementById("loadingMsg");
const vid = document.getElementById("vid");

// Force inline playback on iOS — the JS property (camelCase) is what
// iOS actually checks; the HTML attribute alone is not always enough,
// especially in WebViews (Mail, Files, Messages).
vid.playsInline = true;
vid.setAttribute("playsinline", "");
vid.setAttribute("webkit-playsinline", "");

try {{
    const byteChars = atob(videoB64);
    const len = byteChars.length;
    const bytes = new Uint8Array(len);
    for (let i = 0; i < len; i++) bytes[i] = byteChars.charCodeAt(i);
    const blob = new Blob([bytes], {{ type: "video/mp4" }});
    const blobUrl = URL.createObjectURL(blob);
    vid.src = blobUrl;
    if (loadMsg) loadMsg.style.display = "none";
}} catch(e) {{
    if (loadMsg) loadMsg.textContent = "Error loading video: " + e.message;
    console.error("Video blob creation failed:", e);
}}
const segs = document.querySelectorAll(".seg");
const panel = document.getElementById("scriptPanel");
const notesBar = document.getElementById("notesBar");
const notesText = document.getElementById("notesText");
const notesData = [{notes_js_array}];

// Speed control
document.querySelectorAll(".speed-btn").forEach(btn => {{
    btn.addEventListener("click", () => {{
        const speed = parseFloat(btn.dataset.speed);
        vid.playbackRate = speed;
        document.querySelectorAll(".speed-btn").forEach(b => b.classList.remove("active"));
        btn.classList.add("active");
    }});
}});

segs.forEach(s => {{
    s.addEventListener("click", () => {{
        vid.currentTime = parseFloat(s.dataset.start);
        vid.play();
    }});
}});

let lastNotesSlide = -1;
vid.addEventListener("timeupdate", () => {{
    const t = vid.currentTime;
    let activeEl = null;
    segs.forEach(s => {{
        const start = parseFloat(s.dataset.start);
        const end = parseFloat(s.dataset.end);
        if (t >= start && t < end) {{
            s.classList.add("active");
            activeEl = s;
        }} else {{
            s.classList.remove("active");
        }}
    }});
    if (activeEl) {{
        const panelRect = panel.getBoundingClientRect();
        const elRect = activeEl.getBoundingClientRect();
        const offset = elRect.top - panelRect.top - panelRect.height / 3;
        if (Math.abs(offset) > 20) {{
            panel.scrollBy({{ top: offset, behavior: "smooth" }});
        }}
    }}
    // Update notes bar below video
    if (notesBar && notesText) {{
        const nd = notesData.find(n => t >= n.start && t < n.end);
        const slideNum = nd ? nd.slide : -1;
        if (slideNum !== lastNotesSlide) {{
            lastNotesSlide = slideNum;
            if (nd && nd.notes) {{
                notesBar.classList.remove("empty");
                notesText.classList.remove("placeholder");
                notesText.innerHTML = "<strong>Slide " + nd.slide + ":</strong> " + nd.notes;
            }} else {{
                notesBar.classList.add("empty");
                notesText.classList.add("placeholder");
                notesText.innerHTML = nd ? "No notes for slide " + nd.slide : "Notes will appear here as the presentation plays";
            }}
        }}
    }}
}});
</script>
</body>
</html>'''

    with open(viewer_path, "w", encoding="utf-8") as f:
        f.write(html_content)

    size_mb = os.path.getsize(viewer_path) / (1024 * 1024)
    print(f"      Standalone viewer saved to {viewer_path} ({size_mb:.1f} MB)")
    return viewer_path


def generate_mobile_viewer(video_path: str, timeline: list) -> str:
    """
    Generate a mobile-friendly HTML viewer that references the video as a
    separate file instead of embedding it as base64.  This keeps the HTML
    tiny (<50 KB) so mobile Safari can load it, and the video streams
    normally with playsinline support.

    Both the HTML and a copy of the MP4 are placed in a *_mobile/ folder
    next to the original video so they can be uploaded together (e.g. to
    Dropbox, Google Drive, iCloud, etc.).

    Returns the path to the generated HTML file.
    """
    import html as html_mod

    if not timeline or not os.path.exists(video_path):
        return ""

    video_name = os.path.splitext(os.path.basename(video_path))[0]
    mobile_dir = os.path.splitext(video_path)[0] + "_mobile"
    os.makedirs(mobile_dir, exist_ok=True)

    # Copy the MP4 into the mobile folder
    mp4_filename = os.path.basename(video_path)
    mobile_mp4 = os.path.join(mobile_dir, mp4_filename)
    if not os.path.exists(mobile_mp4) or (
        os.path.getsize(mobile_mp4) != os.path.getsize(video_path)
    ):
        shutil.copy2(video_path, mobile_mp4)

    viewer_html_path = os.path.join(mobile_dir, f"{video_name}_viewer.html")

    total_duration = timeline[-1]["end"] if timeline else 0
    total_min = int(total_duration // 60)
    total_sec = int(total_duration % 60)

    # Check if any slides have notes
    has_any_notes = any(t.get("notes", "") for t in timeline)

    # Build notes bar HTML
    if has_any_notes:
        notes_bar_html = (
            '<div class="notes-bar empty" id="notesBar">'
            '<div class="notes-bar-label">Slide Notes</div>'
            '<div class="notes-bar-text placeholder" id="notesText">'
            'Notes will appear here as the presentation plays</div></div>'
        )
    else:
        notes_bar_html = ""

    # Build segment HTML blocks and notes JS data
    seg_blocks = ""
    notes_js_entries = []
    for i, t in enumerate(timeline):
        escaped_text = html_mod.escape(t["text"]).replace("\n", "<br>")
        start_min = int(t["start"] // 60)
        start_sec = int(t["start"] % 60)
        notes_html = ""
        notes_text = t.get("notes", "")
        if notes_text:
            escaped_notes = html_mod.escape(notes_text).replace("\n", "<br>")
            notes_html = (
                f'<div class="seg-notes">'
                f'<span class="seg-notes-label">Slide Notes</span>'
                f'{escaped_notes}</div>'
            )
        js_notes = html_mod.escape(notes_text).replace("\\", "\\\\").replace('"', '\\"').replace("\n", "<br>")
        notes_js_entries.append(f'{{start:{t["start"]},end:{t["end"]},slide:{t["slide"]},notes:"{js_notes}"}}')
        seg_blocks += f'''
        <div class="seg" id="seg-{i}" data-start="{t['start']}" data-end="{t['end']}">
            <div class="seg-header">
                Slide {t['slide']}
                <span class="seg-time">{start_min}:{start_sec:02d}</span>
            </div>
            <div class="seg-text">{escaped_text}</div>
            {notes_html}
        </div>'''
    notes_js_array = ",".join(notes_js_entries)

    html_content = f'''<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>Presentation Viewer — {len(timeline)} slides, {total_min}m {total_sec}s</title>
<style>
* {{ margin: 0; padding: 0; box-sizing: border-box; }}
body {{ font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
       background: #1a1a2e; color: #e0e0e0; }}
.container {{ display: flex; height: 100vh; }}
.video-panel {{ flex: 1; display: flex; flex-direction: column; align-items: center;
               justify-content: center; padding: 20px; background: #16213e; }}
.video-panel video {{ max-width: 100%; max-height: 70vh; border-radius: 8px;
                      box-shadow: 0 4px 20px rgba(0,0,0,0.5); }}
.video-panel h2 {{ color: #a0c4ff; margin-bottom: 12px; font-size: 14px;
                   letter-spacing: 1px; text-transform: uppercase; }}
.speed-bar {{ display: flex; gap: 6px; margin-top: 14px; align-items: center; }}
.speed-bar span {{ font-size: 12px; color: #7a8ba8; margin-right: 4px; }}
.speed-btn {{ background: rgba(255,255,255,0.08); border: 1px solid #334155;
              color: #94a3b8; padding: 5px 12px; border-radius: 6px; cursor: pointer;
              font-size: 13px; transition: all 0.15s; }}
.speed-btn:hover {{ background: rgba(255,255,255,0.15); color: #e0e0e0; }}
.speed-btn.active {{ background: #539cff; color: #fff; border-color: #539cff; }}
.script-panel {{ width: 420px; min-width: 350px; overflow-y: auto; padding: 20px;
                background: #0f3460; border-left: 2px solid #1a1a2e; }}
.script-panel h2 {{ color: #a0c4ff; margin-bottom: 16px; font-size: 14px;
                    letter-spacing: 1px; text-transform: uppercase;
                    position: sticky; top: 0; background: #0f3460; padding: 8px 0; z-index: 1; }}
.seg {{ padding: 14px 16px; margin-bottom: 10px; border-radius: 8px;
        border-left: 3px solid transparent; transition: all 0.3s ease;
        cursor: pointer; background: rgba(255,255,255,0.03); }}
.seg:hover {{ background: rgba(255,255,255,0.08); }}
.seg.active {{ background: rgba(83,156,255,0.15); border-left-color: #539cff; }}
.seg-header {{ font-weight: 600; color: #539cff; margin-bottom: 6px;
              display: flex; justify-content: space-between; align-items: center; font-size: 13px; }}
.seg-time {{ font-weight: 400; color: #7a8ba8; font-size: 12px; }}
.seg-text {{ font-size: 14px; line-height: 1.6; color: #c8d6e5; }}
.seg.active .seg-text {{ color: #f0f0f0; }}
.seg-notes {{ margin-top: 10px; padding: 10px 12px; background: rgba(251,191,36,0.12);
              border: 1px solid rgba(251,191,36,0.3); border-radius: 6px;
              font-size: 13px; line-height: 1.55; color: #fde68a; }}
.seg-notes-label {{ display: block; font-size: 11px; font-weight: 600;
                    text-transform: uppercase; letter-spacing: 0.5px;
                    color: #fbbf24; margin-bottom: 4px; }}
.notes-bar {{ width: 100%; margin-top: 14px; padding: 12px 16px;
             background: rgba(251,191,36,0.10); border: 1px solid rgba(251,191,36,0.25);
             border-radius: 8px; min-height: 48px; max-height: 120px; overflow-y: auto;
             transition: opacity 0.3s ease; }}
.notes-bar.empty {{ opacity: 0.4; }}
.notes-bar-label {{ font-size: 11px; font-weight: 600; text-transform: uppercase;
                    letter-spacing: 0.5px; color: #fbbf24; margin-bottom: 4px; }}
.notes-bar-text {{ font-size: 13px; line-height: 1.5; color: #fde68a; }}
.notes-bar-text.placeholder {{ color: #7a8ba8; font-style: italic; }}
.script-panel::-webkit-scrollbar {{ width: 6px; }}
.script-panel::-webkit-scrollbar-track {{ background: transparent; }}
.script-panel::-webkit-scrollbar-thumb {{ background: #475569; border-radius: 3px; }}
@media (max-width: 768px) {{
    body {{ height: auto; overflow-y: auto; }}
    .container {{ flex-direction: column; height: auto; overflow-y: visible; }}
    .video-panel {{ padding: 12px; flex: none; }}
    .video-panel video {{ max-height: 35vh; width: 100%; }}
    .script-panel {{ width: 100%; min-width: unset; max-height: none;
                     overflow-y: visible;
                     border-left: none; border-top: 2px solid #1a1a2e; }}
    .speed-bar {{ flex-wrap: wrap; }}
    .notes-bar {{ max-height: 100px; }}
}}
</style>
</head>
<body>
<div class="container">
    <div class="video-panel">
        <h2>Presentation</h2>
        <video id="vid" controls playsinline webkit-playsinline>
            <source src="{mp4_filename}" type="video/mp4">
        </video>
        <div class="speed-bar">
            <span>Speed:</span>
            <button class="speed-btn" data-speed="0.5">0.5x</button>
            <button class="speed-btn" data-speed="0.75">0.75x</button>
            <button class="speed-btn active" data-speed="1">1x</button>
            <button class="speed-btn" data-speed="1.25">1.25x</button>
            <button class="speed-btn" data-speed="1.5">1.5x</button>
            <button class="speed-btn" data-speed="2">2x</button>
        </div>
        {notes_bar_html}
    </div>
    <div class="script-panel" id="scriptPanel">
        <h2>Speaker Script</h2>
{seg_blocks}
    </div>
</div>
<script>
const vid = document.getElementById("vid");
vid.playsInline = true;
const segs = document.querySelectorAll(".seg");
const panel = document.getElementById("scriptPanel");
const notesBar = document.getElementById("notesBar");
const notesText = document.getElementById("notesText");
const notesData = [{notes_js_array}];

// Speed control
document.querySelectorAll(".speed-btn").forEach(btn => {{
    btn.addEventListener("click", () => {{
        const speed = parseFloat(btn.dataset.speed);
        vid.playbackRate = speed;
        document.querySelectorAll(".speed-btn").forEach(b => b.classList.remove("active"));
        btn.classList.add("active");
    }});
}});

segs.forEach(s => {{
    s.addEventListener("click", () => {{
        vid.currentTime = parseFloat(s.dataset.start);
        vid.play();
    }});
}});

let lastNotesSlide = -1;
vid.addEventListener("timeupdate", () => {{
    const t = vid.currentTime;
    let activeEl = null;
    segs.forEach(s => {{
        const start = parseFloat(s.dataset.start);
        const end = parseFloat(s.dataset.end);
        if (t >= start && t < end) {{
            s.classList.add("active");
            activeEl = s;
        }} else {{
            s.classList.remove("active");
        }}
    }});
    if (activeEl) {{
        const panelRect = panel.getBoundingClientRect();
        const elRect = activeEl.getBoundingClientRect();
        const offset = elRect.top - panelRect.top - panelRect.height / 3;
        if (Math.abs(offset) > 20) {{
            panel.scrollBy({{ top: offset, behavior: "smooth" }});
        }}
    }}
    if (notesBar && notesText) {{
        const nd = notesData.find(n => t >= n.start && t < n.end);
        const slideNum = nd ? nd.slide : -1;
        if (slideNum !== lastNotesSlide) {{
            lastNotesSlide = slideNum;
            if (nd && nd.notes) {{
                notesBar.classList.remove("empty");
                notesText.classList.remove("placeholder");
                notesText.innerHTML = "<strong>Slide " + nd.slide + ":</strong> " + nd.notes;
            }} else {{
                notesBar.classList.add("empty");
                notesText.classList.add("placeholder");
                notesText.innerHTML = nd ? "No notes for slide " + nd.slide : "Notes will appear here as the presentation plays";
            }}
        }}
    }}
}});
</script>
</body>
</html>'''

    with open(viewer_html_path, "w", encoding="utf-8") as f:
        f.write(html_content)

    html_kb = os.path.getsize(viewer_html_path) / 1024
    mp4_mb = os.path.getsize(mobile_mp4) / (1024 * 1024)
    print(f"      Mobile viewer folder: {mobile_dir}/")
    print(f"        HTML: {os.path.basename(viewer_html_path)} ({html_kb:.0f} KB)")
    print(f"        Video: {mp4_filename} ({mp4_mb:.1f} MB)")
    return viewer_html_path
